#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""이음 IEUM — GASOK 피치덱 생성기 (16:9, 다크+골드, 디자이너 그레이드)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
LOGO = os.path.join(ROOT, "web", "public", "logo.png")
OUT = os.path.join(HERE, "IEUM_Pitch_Deck.pptx")

# ── 팔레트 ───────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x0A, 0x0B)
BG2     = RGBColor(0x10, 0x10, 0x13)
PANEL   = RGBColor(0x16, 0x16, 0x1A)
PANEL2  = RGBColor(0x1C, 0x1C, 0x21)
LINE    = RGBColor(0x2C, 0x2C, 0x33)
INK     = RGBColor(0xF5, 0xF5, 0xF3)
MUTE    = RGBColor(0x9C, 0x9C, 0xA3)
FAINT   = RGBColor(0x66, 0x66, 0x6E)
GOLD    = RGBColor(0xE7, 0xB7, 0x5F)
GOLDSFT = RGBColor(0xF2, 0xCE, 0x8A)
EMER    = RGBColor(0x53, 0xC6, 0x8B)
AMBER   = RGBColor(0xE3, 0xB3, 0x41)
SKY     = RGBColor(0x7A, 0xC6, 0xF0)

DISP = "Segoe UI"        # 라틴 디스플레이
EA   = "Malgun Gothic"   # 한글

EMU_IN = 914400
PRS = Presentation()
PRS.slide_width  = Inches(13.333)
PRS.slide_height = Inches(7.5)
BLANK = PRS.slide_layouts[6]
SW, SH = 13.333, 7.5
MX = 0.92  # 좌우 여백


# ── 헬퍼 ─────────────────────────────────────────────────
def _set_font(run, name=DISP, ea=EA):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea if tag == "a:ea" else name)


def slide():
    s = PRS.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PRS.slide_width, PRS.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, tracking=None, name=DISP, italic=False,
         space_after=0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, color, bold, size)]
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    if space_after:
        p.space_after = Pt(space_after)
    for seg in runs:
        t, c, b, sz = (seg + (None,) * 4)[:4]
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz if sz else size)
        r.font.bold = bool(b) if b is not None else bold
        r.font.italic = italic
        r.font.color.rgb = c if c is not None else color
        _set_font(r, name)
        if tracking is not None:
            r._r.get_or_add_rPr().set("spc", str(int(tracking)))
    return tb


def para(tb, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=1.0,
         tracking=None, name=DISP, space_before=0, space_after=0):
    p = tb.text_frame.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, color, bold, size)]
    for seg in runs:
        t, c, b, sz = (seg + (None,) * 4)[:4]
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz if sz else size)
        r.font.bold = bool(b) if b is not None else bold
        r.font.color.rgb = c if c is not None else color
        _set_font(r, name)
        if tracking is not None:
            r._r.get_or_add_rPr().set("spc", str(int(tracking)))
    return p


def kicker(s, idx, label):
    text(s, MX, 0.62, 8, 0.3, [("이음 IEUM", GOLD, True, 12)], tracking=180)
    text(s, MX, 0.98, 9, 0.35, [(label.upper(), MUTE, True, 12.5)], tracking=260)
    rect(s, MX, 1.36, 0.62, 0.028, fill=GOLD)
    # page number (우상단)
    text(s, SW - MX - 2, 0.62, 2, 0.3, [(f"{idx:02d}", INK, True, 12), (" / 10", FAINT, False, 12)],
         align=PP_ALIGN.RIGHT, tracking=120)


def footer(s):
    text(s, MX, SH - 0.52, 8, 0.3, [("목돈의 길을, 잇다", FAINT, False, 10.5)], tracking=60)
    text(s, SW - MX - 4, SH - 0.52, 4, 0.3, [("ieum-protocol.vercel.app", FAINT, False, 10.5)],
         align=PP_ALIGN.RIGHT, tracking=40)


def headline(s, y, lines, size=40, lead=1.06):
    tb = s.shapes.add_textbox(Inches(MX), Inches(y), Inches(SW - 2 * MX), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = lead
        for seg in ln:
            t, c, b = (seg + (None, None))[:3]
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = True if b is None else b
            r.font.color.rgb = c if c else INK
            _set_font(r, DISP)
    return tb


def chip_num(s, x, y, n, d=0.42, color=GOLD):
    c = rect(s, x, y, d, d, fill=None, line=color, lw=1.4, shape=MSO_SHAPE.OVAL)
    tf = c.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = color; _set_font(r, DISP)
    return c


def card(s, x, y, w, h, fill=PANEL, line=LINE, radius=0.06):
    return rect(s, x, y, w, h, fill=fill, line=line, lw=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def arrow(s, x1, y1, x2, y2, color=GOLD, w=1.6):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    end = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(end)
    return cn


def line(s, x1, y1, x2, y2, color=GOLD, w=1.6):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    return cn


def _sh(s, shape, x, y, w, h, line_c=GOLD, lw=1.7, fill=None):
    return rect(s, x, y, w, h, fill=fill, line=line_c, lw=lw, shape=shape)


def icon(s, kind, cx, cy, r=0.24, color=GOLD, lw=1.9):
    """다크+골드 라인 아이콘 세트 (진짜 벡터 도형)."""
    T = MSO_SHAPE
    if kind == "house":
        _sh(s, T.ISOSCELES_TRIANGLE, cx - r, cy - r, 2 * r, 0.95 * r, color, lw)
        _sh(s, T.RECTANGLE, cx - 0.72 * r, cy - 0.10 * r, 1.44 * r, 1.02 * r, color, lw)
        _sh(s, T.RECTANGLE, cx - 0.2 * r, cy + 0.36 * r, 0.4 * r, 0.56 * r, color, lw)
    elif kind == "box":
        _sh(s, T.CUBE, cx - r, cy - 0.95 * r, 2 * r, 1.9 * r, color, lw)
    elif kind == "pot":
        _sh(s, T.CAN, cx - 0.82 * r, cy - r, 1.64 * r, 2 * r, color, lw)
    elif kind == "money":
        _sh(s, T.ROUNDED_RECTANGLE, cx - r, cy - 0.6 * r, 2 * r, 1.2 * r, color, lw)
        _sh(s, T.OVAL, cx - 0.3 * r, cy - 0.3 * r, 0.6 * r, 0.6 * r, color, lw)
    elif kind == "id":
        _sh(s, T.ROUNDED_RECTANGLE, cx - r, cy - 0.66 * r, 2 * r, 1.32 * r, color, lw)
        _sh(s, T.OVAL, cx - 0.74 * r, cy - 0.34 * r, 0.6 * r, 0.6 * r, color, lw)
        line(s, cx + 0.04 * r, cy - 0.2 * r, cx + 0.74 * r, cy - 0.2 * r, color, 1.4)
        line(s, cx + 0.04 * r, cy + 0.14 * r, cx + 0.74 * r, cy + 0.14 * r, color, 1.4)
    elif kind == "phone":
        _sh(s, T.ROUNDED_RECTANGLE, cx - 0.6 * r, cy - r, 1.2 * r, 2 * r, color, lw)
        line(s, cx - 0.16 * r, cy + 0.66 * r, cx + 0.16 * r, cy + 0.66 * r, color, 1.6)
    elif kind == "bolt":
        _sh(s, T.LIGHTNING_BOLT, cx - 0.62 * r, cy - r, 1.24 * r, 2 * r, line_c=None, lw=lw, fill=color)


# ═══════════════════════════════════════════════════════════
# Slide 1 — 표지
# ═══════════════════════════════════════════════════════════
s = slide()
# 은은한 상단 골드 헤어라인
rect(s, 0, 0, SW, 0.06, fill=GOLD)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(MX), Inches(0.95), height=Inches(1.15))
text(s, MX, 2.5, 10, 1.4, [("이음", INK, True, 88), ("  IEUM", GOLD, True, 58)])
text(s, MX, 3.95, 11, 0.7, [("돈이 사람 손을 거치지 않는, 한국형 온체인 에스크로", INK, False, 24)])
text(s, MX, 4.55, 11, 0.6, [("전세 보증금 · 이사 잔금 · 계모임을 ", MUTE, False, 17),
                            ("단 한 트랜잭션", GOLDSFT, True, 17),
                            ("으로 정산합니다.", MUTE, False, 17)])
rect(s, MX, 5.55, SW - 2 * MX, 0.02, fill=LINE)
text(s, MX, 5.75, 12, 0.4, [("GIWA SEPOLIA", GOLD, True, 12.5), ("      UPBIT × GIWA · GASOK BUILDER PROGRAM",
     MUTE, True, 12.5)], tracking=200)
text(s, SW - MX - 4.5, 5.72, 4.5, 0.4, [("ieum-protocol.vercel.app", INK, True, 13)],
     align=PP_ALIGN.RIGHT, tracking=40)

# ═══════════════════════════════════════════════════════════
# Slide 2 — 문제
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 2, "The Problem"); footer(s)
headline(s, 1.7, [[("한국인의 목돈은 ", INK), ("'사람 손'", GOLD), ("을", INK)],
                  [("거칠 때 터진다", INK)]], size=42)
rows = [("house", "전세금", "집주인을 거치며", "전세사기", AMBER),
        ("box", "이사 잔금", "날짜 하루만 어긋나도", "거래 파탄", AMBER),
        ("pot", "곗돈", "계주 먹튀로", "증발", AMBER)]
y = 3.55
for ic, a, mid, bad, col in rows:
    card(s, MX, y, SW - 2 * MX, 0.82)
    icon(s, ic, MX + 0.72, y + 0.41, r=0.2, color=GOLD)
    text(s, MX + 1.25, y, 3.2, 0.82, [(a, INK, True, 21)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX + 4.3, y, 4.2, 0.82, [(mid, MUTE, False, 16)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW - MX - 2.9, y, 2.5, 0.82, [("→  ", FAINT, False, 16), (bad, col, True, 19)],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    y += 0.9
text(s, MX, 6.38, SW - 2 * MX, 0.4,
     [("셋 다 원인은 하나 — ", MUTE, False, 17), ("목돈이 중개자의 손에 머무는 시간.", GOLDSFT, True, 17)])

# ═══════════════════════════════════════════════════════════
# Slide 3 — 해결
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 3, "The Solution"); footer(s)
headline(s, 1.7, [[("중개자를 ", INK), ("스마트 컨트랙트", GOLD), ("로", INK)],
                  [("대체한다", INK)]], size=42)
tb3 = text(s, MX, 3.65, SW - 2 * MX, 1.2,
     [("신규 세입자 전세금을 ", INK, False, 20), ("온체인에 락", GOLDSFT, True, 20),
      ("한 뒤, 정산일에 단 한 번의 트랜잭션으로", INK, False, 20)], spacing=1.3)
para(tb3, [("기존 세입자 보증금 반환 + 집주인 잔금 지급을 ", INK, False, 20),
           ("동시 확정", GOLDSFT, True, 20), ("합니다.", INK, False, 20)], spacing=1.3, space_before=6)
# 큰 강조 카드 — 2단 분할 (좌: 0초, 우: 불가능) + 중앙 디바이더
cardY, cardH = 5.15, 1.4
card(s, MX, cardY, SW - 2 * MX, cardH, fill=BG2, line=GOLD)
midx = MX + (SW - 2 * MX) / 2
rect(s, midx, cardY + 0.28, 0.012, cardH - 0.56, fill=LINE)  # 세로 디바이더
# 좌측 컬럼
text(s, MX + 0.55, cardY + 0.24, midx - MX - 1.0, 0.35, [("돈이 머무는 시간", MUTE, False, 15)])
text(s, MX + 0.55, cardY + 0.6, midx - MX - 1.0, 0.7, [("= 0초", GOLD, True, 42)])
# 우측 컬럼
text(s, midx + 0.55, cardY + 0.24, SW - MX - midx - 1.0, 0.35,
     [("한쪽만 지급되는 중간 상태가", MUTE, False, 15)])
text(s, midx + 0.55, cardY + 0.62, SW - MX - midx - 1.0, 0.6,
     [("원천적으로 불가능", INK, True, 26)])

# ═══════════════════════════════════════════════════════════
# Slide 4 — 작동 원리 다이어그램
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 4, "How it works"); footer(s)
headline(s, 1.6, [[("원자적 연쇄 정산", INK)]], size=38)
text(s, MX, 2.45, SW - 2 * MX, 0.4, [("정산일 이후 누구나 ", MUTE, False, 15),
     ("settle()", GOLDSFT, True, 15), (" 한 번 → 세 정산이 같은 트랜잭션에서 동시 확정", MUTE, False, 15)])

def actor(x, y, w, title, sub, col=INK):
    card(s, x, y, w, 0.92, fill=PANEL, line=LINE)
    text(s, x + 0.28, y + 0.16, w - 0.5, 0.4, [(title, col, True, 16)])
    text(s, x + 0.28, y + 0.52, w - 0.5, 0.3, [(sub, MUTE, False, 11.5)])

cx, cy, cw, ch = 5.0, 3.6, 3.33, 1.5
card(s, cx, cy, cw, ch, fill=BG2, line=GOLD, radius=0.08)
text(s, cx, cy + 0.28, cw, 0.4, [("JeonseEscrow", GOLD, True, 18)], align=PP_ALIGN.CENTER)
text(s, cx, cy + 0.72, cw, 0.4, [("settle() · 1 TX", INK, True, 14)], align=PP_ALIGN.CENTER)
text(s, cx, cy + 1.04, cw, 0.3, [("중간에 돈 쥐는 사람 없음", MUTE, False, 11)], align=PP_ALIGN.CENTER)

actor(0.92, 3.05, 3.2, "신규 세입자 B", "전세금을 락한다")
actor(0.92, 4.55, 3.2, "이음 Earn", "날짜 공백 · 역전세 선지급")
actor(SW - 0.92 - 3.2, 2.95, 3.2, "기존 세입자 A", "보증금 100% 반환")
actor(SW - 0.92 - 3.2, 4.65, 3.2, "집주인 L", "차액만 수령")

arrow(s, 4.12, 3.5, cx, 3.95)                       # B → escrow
arrow(s, 0.92 + 3.2, 4.9, cx, 4.55)                 # Earn → escrow
arrow(s, cx + cw, 3.95, SW - 0.92 - 3.2, 3.5)       # escrow → A
arrow(s, cx + cw, 4.55, SW - 0.92 - 3.2, 4.9)       # escrow → L
text(s, 4.05, 3.12, 1.2, 0.3, [("전세금 락", GOLDSFT, True, 10.5)], align=PP_ALIGN.CENTER)
text(s, 8.05, 3.12, 1.3, 0.3, [("보증금 반환", GOLDSFT, True, 10.5)], align=PP_ALIGN.CENTER)
text(s, 8.05, 4.98, 1.3, 0.3, [("차액", GOLDSFT, True, 10.5)], align=PP_ALIGN.CENTER)
text(s, 4.05, 4.95, 1.4, 0.3, [("선지급", GOLDSFT, True, 10.5)], align=PP_ALIGN.CENTER)
text(s, MX, 6.35, SW - 2 * MX, 0.4,
     [("역전세(반환 > 신규 전세금)면 부족분을 지갑 또는 담보 대출로 조달해 개설합니다.", MUTE, False, 13.5)],
     align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# Slide 5 — 왜 GIWA인가
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 5, "Why GIWA"); footer(s)
headline(s, 1.6, [[("이음은 ", INK), ("GIWA에서만", GOLD), (" 완성된다", INK)]], size=38)
cards = [("money", "KRW 스테이블코인", "mKRW를 실KRW로 교체하는 순간 = 실생활 전세·저축 앱"),
         ("id", "UP.ID", "집주인·세입자 당사자 검증 + 계 완주 온체인 신용이력"),
         ("phone", "GIWA 월렛", "'내 목돈을 지키는 탭'으로 인앱 탑재"),
         ("bolt", "1초 블록", "즉시 정산 UX — 기다림 없는 사용자 경험")]
gw = (SW - 2 * MX - 0.4) / 2
gh = 1.5
positions = [(MX, 2.7), (MX + gw + 0.4, 2.7), (MX, 2.7 + gh + 0.35), (MX + gw + 0.4, 2.7 + gh + 0.35)]
for (ic, t, d), (x, y) in zip(cards, positions):
    card(s, x, y, gw, gh)
    icon(s, ic, x + 0.62, y + 0.62, r=0.22, color=GOLDSFT)
    text(s, x + 1.15, y + 0.28, gw - 1.4, 0.4, [(t, GOLDSFT, True, 18)])
    text(s, x + 1.15, y + 0.72, gw - 1.4, 0.7, [(d, MUTE, False, 13.5)], spacing=1.15)

# ═══════════════════════════════════════════════════════════
# Slide 6 — 제품 패밀리
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 6, "Product Family"); footer(s)
headline(s, 1.6, [[("하나의 에스크로 프리미티브 → ", INK), ("세 제품", GOLD)]], size=34)
items = [("01", "전세 에스크로", "원자적 정산으로 보증금 반환·잔금을 동시 확정", "메인", GOLD),
         ("02", "이음 Earn", "담보 대출 + 전세 브리지가 한 풀에서 도는 통합 머니마켓", "머니마켓", SKY),
         ("03", "계모임", "계주 없는 온체인 계 — 조작 불가 추첨 + 미납 슬래싱", "확장", EMER)]
y = 2.85
for n, t, d, tag, col in items:
    card(s, MX, y, SW - 2 * MX, 1.15)
    chip_num(s, MX + 0.4, y + 0.36, n, color=col)
    text(s, MX + 1.25, y + 0.22, 6, 0.5, [(t, INK, True, 21)])
    text(s, MX + 1.25, y + 0.66, 8.2, 0.4, [(d, MUTE, False, 14)])
    tagbox = rect(s, SW - MX - 1.8, y + 0.38, 1.4, 0.4, fill=None, line=col, lw=1.2,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = tagbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = tag; rr.font.size = Pt(12); rr.font.bold = True
    rr.font.color.rgb = col; _set_font(rr)
    y += 1.33

# ═══════════════════════════════════════════════════════════
# Slide 7 — 시장
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 7, "Market"); footer(s)
headline(s, 1.7, [[("전세는 ", INK), ("한국", GOLD), (", 계는 ", INK), ("전 세계", GOLD)]], size=40)
card(s, MX, 3.5, (SW - 2 * MX - 0.4) / 2, 2.5)
text(s, MX + 0.5, 3.85, 5, 0.5, [("KR", GOLD, True, 13)], tracking=200)
text(s, MX + 0.5, 4.2, 5.5, 0.9, [("수백조 원", INK, True, 46)])
text(s, MX + 0.5, 5.2, 5.5, 0.7, [("규모의 전세 보증금이 매년 이동.", MUTE, False, 15),
     ("\n한 건당 수천만 원의 TVL이 온체인으로.", MUTE, False, 15)], spacing=1.2)
x2 = MX + (SW - 2 * MX - 0.4) / 2 + 0.4
card(s, x2, 3.5, (SW - 2 * MX - 0.4) / 2, 2.5)
text(s, x2 + 0.5, 3.85, 5, 0.5, [("GLOBAL", GOLD, True, 13)], tracking=200)
text(s, x2 + 0.5, 4.2, 5.5, 0.9, [("ROSCA", INK, True, 46)])
text(s, x2 + 0.5, 5.2, 5.5, 0.7, [("계는 전 세계 현상 — 동남아 ", MUTE, False, 15),
     ("arisan", INK, True, 15), (", 중남미 ", MUTE, False, 15), ("tanda", INK, True, 15),
     (".\n수출 가능한 프리미티브.", MUTE, False, 15)], spacing=1.2)

# ═══════════════════════════════════════════════════════════
# Slide 8 — 신뢰 & 보안
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 8, "Trust & Security"); footer(s)
headline(s, 1.55, [[("실자금을 다루는 팀의 태도를", INK)], [("코드로 증명한다", INK)]], size=34, lead=1.05)
# 좌: 지표
stats = [("68", "Foundry 테스트 전부 통과", EMER),
         ("3", "자가 감사로 치명 취약점 발견·수정", GOLD),
         ("6", "컨트랙트 GIWA 익스플로러 Verified", SKY)]
y = 3.5
for big, lab, col in stats:
    text(s, MX, y, 1.6, 0.7, [(big, col, True, 38)])
    text(s, MX + 1.6, y + 0.12, 4.2, 0.7, [(lab, MUTE, False, 14.5)], anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    y += 0.98
# 우: 방어선 카드
rx = SW / 2 + 0.2
card(s, rx, 3.4, SW - MX - rx, 3.05, fill=PANEL, line=LINE)
text(s, rx + 0.45, 3.68, 5, 0.4, [("내장된 방어선", GOLDSFT, True, 15)], tracking=40)
defs = [("위조 에스크로 드레인", "팩토리 화이트리스트로 차단"),
        ("추첨 그라인딩", "커밋–리빌 2단계 온체인 추첨"),
        ("오라클 급변 / 데스 스파이럴", "서킷브레이커 ±20% + 긴급정지"),
        ("멈춘 오라클", "스테일 가드로 과다대출·부당청산 방지")]
yy = 4.2
for a, b in defs:
    text(s, rx + 0.45, yy, 0.4, 0.4, [("✓", EMER, True, 15)])
    text(s, rx + 0.85, yy - 0.02, SW - MX - rx - 1.1, 0.6,
         [(a, INK, True, 13.5), ("  ·  ", FAINT, False, 12), (b, MUTE, False, 12.5)], spacing=1.0)
    yy += 0.56

# ═══════════════════════════════════════════════════════════
# Slide 9 — 팀 & 실행력
# ═══════════════════════════════════════════════════════════
s = slide(); kicker(s, 9, "Team & Execution"); footer(s)
headline(s, 1.7, [[("1인 + ", INK), ("AI 네이티브 실행력", GOLD)]], size=40)
text(s, MX, 3.3, SW - 2 * MX, 0.7,
     [("아이디어 확정 → GIWA 테스트넷 배포·실사용 완주까지 ", INK, False, 19),
      ("【N】일", GOLDSFT, True, 19), (".", INK, False, 19)])
metrics = [("MVP", "동작하는 라이브 dApp"), ("6", "Verified 컨트랙트"),
           ("68", "통과 테스트"), ("3", "제품 (에스크로·Earn·계)")]
gw = (SW - 2 * MX - 3 * 0.3) / 4
for i, (big, lab) in enumerate(metrics):
    x = MX + i * (gw + 0.3)
    card(s, x, 4.35, gw, 1.5)
    text(s, x, 4.6, gw, 0.7, [(big, GOLD, True, 34)], align=PP_ALIGN.CENTER)
    text(s, x + 0.15, 5.42, gw - 0.3, 0.4, [(lab, MUTE, False, 12.5)], align=PP_ALIGN.CENTER, spacing=1.05)
text(s, MX, 6.25, SW - 2 * MX, 0.4, [("【실명 · 역할 · 한 줄 이력】 — 말보다 배포된 컨트랙트로 증명합니다.",
     FAINT, False, 13)])

# ═══════════════════════════════════════════════════════════
# Slide 10 — Ask / 마무리
# ═══════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, 0.06, fill=GOLD)
kicker(s, 10, "The Ask")
headline(s, 1.7, [[("GASOK에 요청드립니다", INK)]], size=40)
asks = ["스마트 컨트랙트 외부 감사",
        "KRW 스테이블코인 조기 테스트 액세스",
        "UP.ID 연동 기술 지원",
        "GIWA 월렛 인앱 탑재 검토"]
y = 3.15
for i, a in enumerate(asks, 1):
    chip_num(s, MX, y, i)
    text(s, MX + 0.7, y - 0.03, 9, 0.5, [(a, INK, True, 20)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.72
rect(s, MX, 6.15, SW - 2 * MX, 0.02, fill=LINE)
text(s, MX, 6.35, 9, 0.5, [("이음 · IEUM", INK, True, 20), ("  —  목돈의 길을, 잇다.", GOLD, True, 20)])
text(s, SW - MX - 5, 6.4, 5, 0.4, [("ieum-protocol.vercel.app", GOLDSFT, True, 15)],
     align=PP_ALIGN.RIGHT)

PRS.save(OUT)
print("saved:", OUT, "| slides:", len(PRS.slides._sldIdLst))
